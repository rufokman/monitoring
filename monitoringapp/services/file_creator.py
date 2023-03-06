import datetime
import pandas as pd
from django.utils import timezone
from django.shortcuts import get_object_or_404
import xlwt
from django.http import HttpResponse
from ..models import *
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import io


def get_update_data_not_fix():

	card_data = Card.objects.filter(delete=0)

	return card_data


def download_excel_admin(request):
	today = "KPI_report_{}".format(datetime.datetime.today().strftime("%d.%m.%Y"))
	response = HttpResponse(content_type='application/ms-excel')
	response['Content-Disposition'] = 'attachment; filename={}.xlsx'.format(today)
	workbook = xlsxwriter.Workbook(response, {'in_memory': True})
	format = workbook.add_format()
	format.set_bg_color('#fdc433')
	ws = workbook.add_worksheet('Реестр')
	format_date = workbook.add_format({'num_format': 'dd.mm.yyyy hh:mm:ss'})

	ws.set_column('A:A', 5)
	ws.set_column('B:B', 23)
	ws.set_column('C:C', 23, )
	ws.set_column('D:D', 23, )
	ws.set_column('E:E', 26, )
	ws.set_column('F:F', 11, )
	ws.set_column('G:G', 81, )
	ws.set_column('H:H', 16, )
	ws.set_column('I:I', 20, )
	ws.set_column('J:J', 20, )
	ws.set_column('K:K', 30, )
	ws.set_column('L:L', 20, )
	ws.set_column('M:M', 7, )
	ws.set_column('N:N', 17, )
	ws.set_column('O:O', 17, )
	ws.set_column('P:P', 40, )
	ws.set_column('Q:Q', 24, )
	ws.set_column('R:R', 24, )
	ws.set_column('S:S', 24, )
	columns = ['№ пп',  'Организация', 'ФИО сотрудника, в чью карту устанавливается КПЭ', 'Название должности',
			   'КПЭ / КлС2', 'Наименование КПЭ / КлС', 'Тип КПЭ/КлС',
			    'Нижний уровень', "Целевой уровень",
			   "Верхний уровень", "Вес КПЭ/Значимость КлС", "Комментарий СУП УК",
			   "Факт 1 кв", "Прогноз 1 кв", "Факт 2 кв", "Прогноз 2 кв", "Факт 3 кв", "Прогноз 3 кв",
			   "Причина отклонения", "Мероприятия по снижению риска невыполнения", "Прогнозируемый результат после реализации мероприятий",
			   "Верификатор", "Последнее обновление", "ФИО пользователя", "Статуc"
			   ]
	row_num = 0
	for col_num in range(len(columns)):
		ws.write(row_num, col_num, columns[col_num])

	row_num=0
	for my_row in get_update_data_not_fix():
		row_num = row_num + 1
		ws.write(row_num, 0, row_num)
		ws.write(row_num, 1, my_row.organization)
		ws.write(row_num, 2, my_row.fio)
		ws.write(row_num, 3, my_row.role)
		ws.write(row_num, 4, my_row.type)
		ws.write(row_num, 5, my_row.name)
		ws.write(row_num, 6, my_row.method)
		ws.write(row_num, 7, my_row.low_level)
		ws.write(row_num, 8, my_row.target_level)
		ws.write(row_num, 9, my_row.high_level)
		ws.write(row_num, 10, my_row.weight)
		ws.write(row_num, 11, my_row.comment)
		ws.write(row_num, 12, my_row.first_quarter_fact)
		ws.write(row_num, 13, my_row.first_quarter)
		ws.write(row_num, 14,  my_row.second_quarter_fact)
		ws.write(row_num, 15,  my_row.second_quarter)
		ws.write(row_num, 16, my_row.third_quarter_fact)
		ws.write(row_num, 17, my_row.third_quarter)
		ws.write(row_num, 18, my_row.reason)
		ws.write(row_num, 19, my_row.measure)
		ws.write(row_num, 20, my_row.forecast)
		ws.write(row_num, 21, my_row.verificator)
		ws.write(row_num, 22, my_row.updated_at.replace(tzinfo=None), format_date)
		ws.write(row_num, 23, my_row.name_of_user)
		ws.write(row_num, 24, my_row.status_list[my_row.status][1])

	workbook.close()
	return response


def get_update_data_not_fix_log():

	card_data_log = CardLog.objects.filter(delete=0)

	return card_data_log


def download_excel_admin_log(request):
	today = "Log_{}".format(datetime.datetime.today().strftime("%d.%m.%Y"))
	response = HttpResponse(content_type='application/ms-excel')
	response['Content-Disposition'] = 'attachment; filename={}.xlsx'.format(today)
	workbook = xlsxwriter.Workbook(response, {'in_memory': True})
	format = workbook.add_format()
	format.set_bg_color('#fdc433')
	ws = workbook.add_worksheet('Лог')
	format_date = workbook.add_format({'num_format': 'dd.mm.yyyy hh:mm:ss'})

	ws.set_column('A:A', 5)
	ws.set_column('B:B', 23)
	ws.set_column('C:C', 23, )
	ws.set_column('D:D', 23, )
	ws.set_column('E:E', 26, )
	ws.set_column('F:F', 11, )
	ws.set_column('G:G', 81, )
	ws.set_column('H:H', 16, )
	ws.set_column('I:I', 20, )
	ws.set_column('J:J', 20, )
	ws.set_column('K:K', 30, )
	ws.set_column('L:L', 20, )
	ws.set_column('M:M', 7, )
	ws.set_column('N:N', 17, )
	ws.set_column('O:O', 17, )
	ws.set_column('P:P', 40, )
	ws.set_column('Q:Q', 24, )
	ws.set_column('R:R', 24, )
	ws.set_column('S:S', 24, )
	columns = ['№ пп',  'Организация', 'ФИО сотрудника, в чью карту устанавливается КПЭ', 'Название должности',
			   'КПЭ / КлС2', 'Наименование КПЭ / КлС', 'Тип КПЭ/КлС',
			    'Нижний уровень', "Целевой уровень",
			   "Верхний уровень", "Вес КПЭ/Значимость КлС", "Комментарий СУП УК",
			   "Факт 1 кв", "Прогноз 1 кв", "Факт 2 кв", "Прогноз 2 кв", "Факт 3 кв", "Прогноз 3 кв",
			   "Причина отклонения", "Мероприятия по снижению риска невыполнения", "Прогнозируемый результат после реализации мероприятий",
			   "Верификатор", "Последнее обновление", "ФИО пользователя", "Статуc"]
	row_num = 0
	for col_num in range(len(columns)):
		ws.write(row_num, col_num, columns[col_num])

	row_num = 0
	for my_row in get_update_data_not_fix_log():
		row_num = row_num + 1
		ws.write(row_num, 0, row_num)
		ws.write(row_num, 1, my_row.organization)
		ws.write(row_num, 2, my_row.fio)
		ws.write(row_num, 3, my_row.role)
		ws.write(row_num, 4, my_row.type)
		ws.write(row_num, 5, my_row.name)
		ws.write(row_num, 6, my_row.method)
		ws.write(row_num, 7, my_row.low_level)
		ws.write(row_num, 8, my_row.target_level)
		ws.write(row_num, 9, my_row.high_level)
		ws.write(row_num, 10, my_row.weight)
		ws.write(row_num, 11, my_row.comment)
		ws.write(row_num, 12, my_row.first_quarter_fact)
		ws.write(row_num, 13, my_row.first_quarter)
		ws.write(row_num, 14,  my_row.second_quarter_fact)
		ws.write(row_num, 15,  my_row.second_quarter)
		ws.write(row_num, 16, my_row.third_quarter_fact)
		ws.write(row_num, 17, my_row.third_quarter)
		ws.write(row_num, 18, my_row.reason)
		ws.write(row_num, 19, my_row.measure)
		ws.write(row_num, 20, my_row.forecast)
		ws.write(row_num, 21, my_row.verificator)
		ws.write(row_num, 22, my_row.updated_at.replace(tzinfo=None), format_date)
		ws.write(row_num, 23, my_row.name_of_user)
		ws.write(row_num, 24, my_row.status_list[my_row.status][1])

	workbook.close()
	return response


def download_admin_pres(request):
	today = "Consolidated_pres_{}".format(datetime.datetime.today().strftime("%d.%m.%Y"))
	prs = Presentation()
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.shapes[0].text = 'Title'

	data = get_update_data_not_fix()
	organization_list = []  # Создаем массив, который заполнится из базы данных
	fio_list = []
	name_list=[]
	for my_row in data:
		organization_list.append(my_row.organization)
		fio_list.append(my_row.fio)
		name_list.append(my_row.name)

	old_counter = 0
	counter = 0
	row_num = 3
	while counter < 21:
	# while counter < len(organization_list):
		if len(organization_list) - counter < row_num:
			counter += len(organization_list) - counter
		else:
			counter += row_num
		col_num = 3
		slide = prs.slides.add_slide(prs.slide_layouts[6])
		x, y, cx, cy = Inches(0.2), Inches(1.3), Inches(9.6), Inches(3)
		shape = slide.shapes.add_table(counter - old_counter, col_num, x, y, cx, cy)
		table = shape.table
		tbl = shape._element.graphic.graphicData.tbl
		style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
		tbl[0][-1].text = style_id
		table.columns[0].width = Inches(2)
		for i in range(old_counter, counter):
			for j in range(col_num):
				cell = table.cell(i - old_counter, j)
				if j % col_num == 0:
					cell.text = organization_list[i]
				if j % col_num == 1:
					cell.text = fio_list[i]
				if j % col_num == 2:
					cell.text = name_list[i]
		old_counter = counter

	response = HttpResponse(content_type='application/vnd.ms-powerpoint')
	response['Content-Disposition'] = 'attachment; filename="{}.pptx"'.format(today)
	source_stream = BytesIO()
	prs.save(source_stream)
	ppt = source_stream.getvalue()
	source_stream.close()
	response.write(ppt)
	return response


def download_excel_user(request):
	today = "KPI_pivot_{}".format(datetime.datetime.today().strftime("%d.%m.%Y"))
	response = HttpResponse(content_type='application/ms-excel')
	response['Content-Disposition'] = 'attachment; filename={}.xlsx'.format(today)
	workbook = xlsxwriter.Workbook(response, {'in_memory': True})
	format = workbook.add_format()
	format.set_bg_color('#fdc433')
	ws = workbook.add_worksheet('Реестр')

	ws.set_column('A:A', 5)
	ws.set_column('B:B', 23)
	ws.set_column('C:C', 23, )
	ws.set_column('D:D', 23, )
	ws.set_column('E:E', 26, )
	ws.set_column('F:F', 11, )
	ws.set_column('G:G', 81, )
	ws.set_column('H:H', 16, )
	ws.set_column('I:I', 20, )
	ws.set_column('J:J', 20, )
	ws.set_column('K:K', 30, )
	ws.set_column('L:L', 20, )
	ws.set_column('M:M', 7, )
	ws.set_column('N:N', 17, )
	ws.set_column('O:O', 17, )
	ws.set_column('P:P', 40, )
	ws.set_column('Q:Q', 24, )
	ws.set_column('R:R', 24, )
	ws.set_column('S:S', 24, )
	columns = ['№ пп', 'Организация', 'ФИО сотрудника, в чью карту устанавливается КПЭ', 'Название должности',
			   'КПЭ / КлС2', 'Наименование КПЭ / КлС', 'Тип КПЭ/КлС',
			    'Нижний уровень', "Целевой уровень",
			   "Верхний уровень", "Вес КПЭ/Значимость КлС", "Комментарий СУП УК",
			   "Факт 1 кв", "Прогноз 1 кв", "Факт 2 кв", "Прогноз 2 кв", "Факт 3 кв", "Прогноз 3 кв",
			   "Причина отклонения", "Мероприятия по снижению риска невыполнения", "Прогнозируемый результат после реализации мероприятий",
			   "Верификатор"
			   ]
	row_num = 0
	for col_num in range(len(columns)):
		ws.write(row_num, col_num, columns[col_num])

	row_num=0
	for my_row in get_update_data_not_fix():
		row_num = row_num + 1
		ws.write(row_num, 0, row_num)
		ws.write(row_num, 1, my_row.organization)
		ws.write(row_num, 2, my_row.fio)
		ws.write(row_num, 3, my_row.role)
		ws.write(row_num, 4, my_row.type)
		ws.write(row_num, 5, my_row.name)
		ws.write(row_num, 6, my_row.method)
		ws.write(row_num, 7, my_row.low_level)
		ws.write(row_num, 8, my_row.target_level)
		ws.write(row_num, 9, my_row.high_level)
		ws.write(row_num, 10, my_row.weight)
		ws.write(row_num, 11, my_row.comment)
		ws.write(row_num, 12, my_row.first_quarter_fact)
		ws.write(row_num, 13, my_row.first_quarter)
		ws.write(row_num, 14,  my_row.second_quarter_fact)
		ws.write(row_num, 15,  my_row.second_quarter)
		ws.write(row_num, 16, my_row.third_quarter_fact)
		ws.write(row_num, 17, my_row.third_quarter)
		ws.write(row_num, 18, my_row.reason)
		ws.write(row_num, 19, my_row.measure)
		ws.write(row_num, 20, my_row.forecast)
		ws.write(row_num, 21, my_row.verificator)

	workbook.close()
	return response